#!/usr/bin/env python3
import sys
import re
import argparse
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_AUTO_SIZE

def add_formatted_text(paragraph, text):
    """Parse basic markdown inline formatting and add as runs to paragraph."""
    # Split by **bold**, __bold__, *italic*, _italic_
    parts = re.split(r'(\*\*.*?\*\*|__.*?__|\*.*?\*|_.*?_)', text)
    for part in parts:
        if not part:
            continue
        run = paragraph.add_run()
        if (part.startswith('**') and part.endswith('**')) or (part.startswith('__') and part.endswith('__')):
            run.text = part[2:-2]
            run.font.bold = True
        elif (part.startswith('*') and part.endswith('*')) or (part.startswith('_') and part.endswith('_')):
            run.text = part[1:-1]
            run.font.italic = True
        else:
            run.text = part

def parse_slide(slide_text, is_first):
    """Parse a single markdown slide block into structured data."""
    data = {
        'layout': 'content',
        'title': '',
        'subtitle': '',
        'bullets': [],
        'images': [],
        'notes': ''
    }
    
    notes_match = re.search(r'<!--\s*notes?:?(.*?)\s*-->', slide_text, re.DOTALL | re.IGNORECASE)
    if notes_match:
        data['notes'] = notes_match.group(1).strip()
        slide_text = slide_text.replace(notes_match.group(0), '')
        
    lines = slide_text.split('\n')
    has_h1 = any(l.strip().startswith('# ') for l in lines)
    if is_first and has_h1:
        data['layout'] = 'title'
        
    for original_line in lines:
        if not original_line.strip():
            continue
        line = original_line.strip()
        
        img_match = re.match(r'!\[.*?\]\((.*?)\)', line)
        if img_match:
            data['images'].append(img_match.group(1))
            continue
            
        if data['layout'] == 'title':
            if line.startswith('# '):
                data['title'] = line[2:].strip()
            elif not line.startswith('#'):
                data['subtitle'] += line + '\n'
        else:
            if line.startswith('## '):
                data['title'] = line[3:].strip()
            elif line.startswith('# '):
                data['title'] = line[2:].strip()
            elif line.startswith('- ') or line.startswith('* '):
                spaces = len(original_line) - len(original_line.lstrip())
                level = spaces // 2
                level = min(level, 8)
                data['bullets'].append({'level': level, 'text': line[2:].strip()})
            else:
                data['bullets'].append({'level': 0, 'text': line})
                
    data['subtitle'] = data['subtitle'].strip()
    return data

def build_presentation(md_file, output_file, template_file=None):
    with open(md_file, 'r', encoding='utf-8') as f:
        md_text = f.read()
        
    raw_slides = re.split(r'\n---\n', md_text)
    
    if template_file and os.path.exists(template_file):
        prs = Presentation(template_file)
    else:
        prs = Presentation()
        # Set to 16:9 Widescreen automatically if no template is provided
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
        
    for i, raw in enumerate(raw_slides):
        raw = raw.strip()
        if not raw:
            continue
            
        slide_data = parse_slide(raw, i == 0)
        
        if slide_data['layout'] == 'title':
            slide_layout = prs.slide_layouts[0] 
            slide = prs.slides.add_slide(slide_layout)
            
            if slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
            if len(slide.placeholders) > 1 and slide_data['subtitle']:
                tf = slide.placeholders[1].text_frame
                tf.word_wrap = True
                tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                tf.clear()
                p = tf.paragraphs[0]
                add_formatted_text(p, slide_data['subtitle'])
                
        else:
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)
            
            if slide.shapes.title:
                slide.shapes.title.text = slide_data['title']
                
            has_text = len(slide_data['bullets']) > 0
            
            if len(slide.placeholders) > 1:
                tf_shape = slide.placeholders[1]
                
                if has_text:
                    tf = tf_shape.text_frame
                    tf.word_wrap = True
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    tf.clear()
                    
                    is_first_p = True
                    for b in slide_data['bullets']:
                        if is_first_p:
                            p = tf.paragraphs[0]
                            is_first_p = False
                        else:
                            p = tf.add_paragraph()
                            
                        p.level = b['level']
                        add_formatted_text(p, b['text'])
                else:
                    tf_shape.text = ""
            
            # --- Image & Layout Logic ---
            if len(slide_data['images']) > 0:
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                img_top = int(slide_height * 0.22)
                img_max_height = int(slide_height * 0.70)
                
                if has_text and len(slide.placeholders) > 1:
                    # Text 40%, Gap 5%, Image 50%, Margin 5%
                    tf_shape.width = int(slide_width * 0.40)
                    
                    img_left = int(slide_width * 0.45)
                    img_max_width = int(slide_width * 0.50)
                else:
                    # Full width image
                    img_left = int(slide_width * 0.1)
                    img_max_width = int(slide_width * 0.8)

                for img_path in slide_data['images']:
                    if os.path.exists(img_path):
                        try:
                            pic = slide.shapes.add_picture(img_path, img_left, img_top)
                            
                            aspect_ratio = pic.width / pic.height
                            target_ratio = img_max_width / img_max_height
                            
                            if aspect_ratio > target_ratio:
                                pic.width = img_max_width
                                pic.height = int(img_max_width / aspect_ratio)
                            else:
                                pic.height = img_max_height
                                pic.width = int(img_max_height * aspect_ratio)
                                
                            pic.left = img_left + int((img_max_width - pic.width) / 2)
                            pic.top = img_top + int((img_max_height - pic.height) / 2)
                            
                        except Exception as e:
                            print(f"Failed to add image {img_path}: {e}")
                    else:
                        print(f"Warning: Image not found: {img_path}")
                    
        if slide_data['notes']:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = slide_data['notes']
            
    prs.save(output_file)
    print(f"✅ Successfully created presentation: {output_file}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert Markdown to PowerPoint PPTX")
    parser.add_argument("input", help="Input Markdown file")
    parser.add_argument("-o", "--output", default="output.pptx", help="Output PPTX file")
    parser.add_argument("-t", "--template", help="Optional PPTX template file")
    
    args = parser.parse_args()
    build_presentation(args.input, args.output, args.template)
